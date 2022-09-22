from codecs import replace_errors
import collections
import collections.abc
from dataclasses import replace
from webbrowser import get

from pptx import Presentation
import pandas as pd
import os

#get cwd
cwd = os.getcwd()

#create dataframe and append xlsx sheets in cwd to it
df = pd.DataFrame()

#concat xlsx sheets in cwd to dataframe
for file in os.listdir(cwd):
    if file.endswith(".xlsx"):
        df = pd.concat([df, pd.read_excel(file)])
    
print(df.size)

#function that replaces a string in a column with another string
def replace_string(df, column, string, replacement):
    df[column] = df[column].str.replace(string, replacement)
    return df

#function that removes all rows with a certain string in a column
def remove_string(df, column, string):
    df = df[df[column] != string]
    return df

#function that removes columns with strings over a certain length
def remove_long_strings(df, column, length):
    df = df[df[column].str.len() <= length]
    return df


#function that returns a dataframe with only the rows that contain a certain string in a column

def get_string(df, column, string):
    df = df[df[column].str.contains(string)]
    return df



#replace emoji characters with blank space
df = replace_string(df,'Ticket satisfaction comment', "_x000D_", " ")

df = get_string(df, 'Ticket satisfaction comment', "Hector")


print(df)

#remove rows with strings over 100 characters
df = remove_long_strings(df, 'Ticket satisfaction comment', 100)

print(df)

#expofrt dataframe to json and save in cwd

#save df as json file to cwd
df.to_json('surveys.json', orient='records')





""""
#open pptx file
prs = Presentation('surveys_master.pptx')

#create slides
for row in df.iterrows():
    prs.slides.add_slide(prs.slide_layouts[6])

#fill slides with text from dataframe
for indx, slide in enumerate(prs.slides):
 
    for shape in slide.placeholders:
      
        if shape.placeholder_format.idx == 18:
            shape.text = df.iloc[indx]['Requester name']
        if shape.placeholder_format.idx == 14:
            shape.text = str(df.iloc[indx]['Ticket ID'])
        if shape.placeholder_format.idx == 16:
            shape.text = df.iloc[indx]['HD System']
        if shape.placeholder_format.idx == 17:
            shape.text = df.iloc[indx]['Solved']
        if shape.placeholder_format.idx == 15:
            shape.text = df.iloc[indx]['Ticket satisfaction comment']
            
            
#save pptx
prs.save(last_name + ' surveys.pptx')
"""